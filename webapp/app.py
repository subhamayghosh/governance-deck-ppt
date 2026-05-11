import os
import json
import io
from datetime import datetime
from flask import Flask, render_template, jsonify, request, send_file
import pandas as pd
from pptx import Presentation
from werkzeug.utils import secure_filename
from slide_updaters import update_presentation

app = Flask(__name__)

from flask.json.provider import DefaultJSONProvider


class SafeJSONProvider(DefaultJSONProvider):
    def default(self, obj):
        if isinstance(obj, pd.Timestamp):
            return None if pd.isna(obj) else obj.isoformat()
        if isinstance(obj, (pd.NaT.__class__,)):
            return None
        if hasattr(obj, 'isoformat'):
            return obj.isoformat()
        if isinstance(obj, float) and (obj != obj):
            return None
        return super().default(obj)


app.json_provider_class = SafeJSONProvider
app.json = SafeJSONProvider(app)

# Upload folder for Excel and PPT files
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Default files (shipped with the app)
DEFAULT_EXCEL = os.path.join(os.path.dirname(os.path.dirname(__file__)), "LPL QE Governance Data_February_2026_MACRO.xlsm")
DEFAULT_PPTX = os.path.join(os.path.dirname(os.path.dirname(__file__)), "2 LPL Cognizant QE Connect Feb 26 V 1.0.pptx")

# State file to persist the currently active files across server restarts
STATE_FILE = os.path.join(UPLOAD_FOLDER, ".active_files.json")


def _load_active_files():
    """Load persisted active files, falling back to defaults if missing/invalid."""
    defaults = {
        "excel": DEFAULT_EXCEL,
        "excel_name": "LPL QE Governance Data_February_2026_MACRO.xlsm",
        "pptx": DEFAULT_PPTX,
        "pptx_name": "2 LPL Cognizant QE Connect Feb 26 V 1.0.pptx",
    }
    if os.path.exists(STATE_FILE):
        try:
            with open(STATE_FILE, "r") as f:
                saved = json.load(f)
            # Verify the referenced files still exist
            if os.path.exists(saved.get("excel", "")) and os.path.exists(saved.get("pptx", "")):
                return saved
        except Exception:
            pass
    return defaults


def _save_active_files():
    """Persist active files state to disk."""
    try:
        with open(STATE_FILE, "w") as f:
            json.dump(active_files, f)
    except Exception as e:
        print(f"Warning: could not persist active files: {e}")


# Track the currently active files (persisted across restarts)
active_files = _load_active_files()


def clean_df(df):
    df.columns = [str(c).strip() for c in df.columns]
    df = df.dropna(how="all").reset_index(drop=True)
    for col in df.columns:
        if pd.api.types.is_datetime64_any_dtype(df[col]):
            df[col] = df[col].dt.strftime("%Y-%m-%d %H:%M:%S").fillna("")
        else:
            df[col] = df[col].fillna("")
    return df


def build_domain_lookup(project_details):
    lookup = {}
    for r in project_details:
        team = str(r.get("Scrum Team Name", "")).strip()
        if team:
            lookup[team.lower()] = (
                str(r.get("Domain", "")),
                str(r.get("Sub Domain", ""))
            )
    return lookup


def fill_missing_domains(records, domain_lookup, team_col="Scrum Team Name"):
    for r in records:
        if not r.get("Domain"):
            team = str(r.get(team_col, "")).strip().lower()
            if team in domain_lookup:
                r["Domain"] = domain_lookup[team][0]
                r["Sub Domain"] = domain_lookup[team][1]
            else:
                r["Domain"] = "Unknown"
                r["Sub Domain"] = ""
    return records


SHEET_MAP = {
    "project_details": "ProjectDetails",
    "tcs_details": "TCsDetailsSheet",
    "insprint_data": "InSprintData",
    "defect_data": "DefectData",
    "monthly_sheet": "MonthlySheet",
    "branch_testcase": "BranchTestCaseSheet",
    "release_day": "ReleaseDayTestCaseSheet",
    "rad_enabled": "RADEnabled",
    "daily_execution": "DailyExecution",
}


def read_excel_data(excel_path=None):
    if excel_path is None:
        excel_path = active_files["excel"]

    data = {}
    xls = pd.ExcelFile(excel_path)
    available_sheets = xls.sheet_names

    for key, sheet_name in SHEET_MAP.items():
        if sheet_name in available_sheets:
            df = pd.read_excel(xls, sheet_name, header=0)
            data[key] = clean_df(df).to_dict(orient="records")
        else:
            data[key] = []

    # Fill missing domains
    if data["project_details"]:
        domain_lookup = build_domain_lookup(data["project_details"])
        data["tcs_details"] = fill_missing_domains(data["tcs_details"], domain_lookup)
        data["branch_testcase"] = fill_missing_domains(data["branch_testcase"], domain_lookup)
        data["insprint_data"] = fill_missing_domains(data["insprint_data"], domain_lookup, team_col="Scrum team")
        data["defect_data"] = fill_missing_domains(data["defect_data"], domain_lookup, team_col="ScrumTeam")
        data["daily_execution"] = fill_missing_domains(data["daily_execution"], domain_lookup)
        data["release_day"] = fill_missing_domains(data["release_day"], domain_lookup)
        data["rad_enabled"] = fill_missing_domains(data["rad_enabled"], domain_lookup)
        data["monthly_sheet"] = fill_missing_domains(data["monthly_sheet"], domain_lookup, team_col="Scrum Team Name")

    return data


def compute_summaries(data):
    summaries = {}

    summaries["total_teams"] = len(data["project_details"])
    active_teams = [r for r in data["project_details"] if str(r.get("Project Status", "")).lower() == "active"]
    summaries["active_teams"] = len(active_teams)

    domains = {}
    for r in data["project_details"]:
        d = str(r.get("Domain", "Unknown"))
        if d:
            domains[d] = domains.get(d, 0) + 1
    summaries["domains"] = domains

    total_tcs = 0
    total_automated = 0
    total_feasible = 0
    total_not_feasible = 0
    for r in data["tcs_details"]:
        for p in ["P0", "P1", "P2", "P3", "P4"]:
            total_tcs += _num(r.get(f"Total TCs {p}", 0))
            total_automated += _num(r.get(f"Total TCs Automated {p}", 0))
            total_feasible += _num(r.get(f"Total TCs Feasible {p}", 0))
            total_not_feasible += _num(r.get(f"Total TCs Not Feasible {p}", 0))
    summaries["total_tcs"] = total_tcs
    summaries["total_automated"] = total_automated
    summaries["total_feasible"] = total_feasible
    summaries["total_not_feasible"] = total_not_feasible
    summaries["automation_coverage"] = round(total_automated / total_feasible * 100, 1) if total_feasible else 0

    defect_by_domain = {}
    for r in data["defect_data"]:
        domain = str(r.get("Domain", "Unknown"))
        severity_total = sum(_num(r.get(k, 0)) for k in [
            "InSprint Fatel", "InSprint Serious", "InSprint Medium", "InSprint Low",
            "Regression Fatel", "Regression Serious", "Regression Medium", "Regression Low"
        ])
        defect_by_domain[domain] = defect_by_domain.get(domain, 0) + severity_total
    summaries["defect_by_domain"] = defect_by_domain

    severity_keys = {
        "Fatal": ["InSprint Fatel", "Regression Fatel"],
        "Serious": ["InSprint Serious", "Regression Serious"],
        "Medium": ["InSprint Medium", "Regression Medium"],
        "Low": ["InSprint Low", "Regression Low"],
    }
    severity_totals = {}
    for label, keys in severity_keys.items():
        severity_totals[label] = sum(_num(r.get(k, 0)) for r in data["defect_data"] for k in keys)
    summaries["defect_severity"] = severity_totals

    total_passed = sum(_num(r.get("Total Passed TC #", 0)) for r in data["daily_execution"])
    total_failed = sum(_num(r.get("Total Failed TC #", 0)) for r in data["daily_execution"])
    summaries["daily_passed"] = total_passed
    summaries["daily_failed"] = total_failed
    summaries["daily_pass_rate"] = round(total_passed / (total_passed + total_failed) * 100, 1) if (total_passed + total_failed) else 0

    rad_by_domain = {}
    for r in data["rad_enabled"]:
        domain = str(r.get("Domain", "Unknown"))
        rad_by_domain.setdefault(domain, {"total_builds": 0, "rad_enabled": 0})
        rad_by_domain[domain]["total_builds"] += _num(r.get("Independent Build #", 0)) + _num(r.get("Shared Build #", 0))
        rad_by_domain[domain]["rad_enabled"] += _num(r.get("# Builds RAD Enabled", 0))
    summaries["rad_by_domain"] = rad_by_domain

    release_by_domain = {}
    for r in data["release_day"]:
        domain = str(r.get("Domain", "Unknown"))
        release_by_domain.setdefault(domain, {"stories": 0, "feasible": 0, "automated": 0})
        release_by_domain[domain]["stories"] += _num(r.get("No Of Stories part of the Release", 0))
        release_by_domain[domain]["feasible"] += _num(r.get("No of Automation feasible TCs", 0))
        release_by_domain[domain]["automated"] += _num(r.get("No of TCs Automated", 0))
    summaries["release_by_domain"] = release_by_domain

    tcs_by_domain = {}
    for r in data["tcs_details"]:
        domain = str(r.get("Domain", "Unknown"))
        total = sum(_num(r.get(f"Total TCs {p}", 0)) for p in ["P0", "P1", "P2", "P3", "P4"])
        auto = sum(_num(r.get(f"Total TCs Automated {p}", 0)) for p in ["P0", "P1", "P2", "P3", "P4"])
        tcs_by_domain.setdefault(domain, {"total": 0, "automated": 0})
        tcs_by_domain[domain]["total"] += total
        tcs_by_domain[domain]["automated"] += auto
    summaries["tcs_by_domain"] = tcs_by_domain

    insprint_teams = {}
    for r in data["insprint_data"]:
        team = str(r.get("Scrum team", ""))
        cov = _num(r.get("Automation coverage", 0))
        if team and cov > 0:
            insprint_teams[team] = round(cov * 100, 1) if cov <= 1 else round(cov, 1)
    summaries["insprint_automation"] = insprint_teams

    return summaries


def _num(v):
    try:
        return float(v) if v != "" else 0
    except (ValueError, TypeError):
        return 0


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/file-info")
def file_info():
    excel_ok = os.path.exists(active_files.get("excel", ""))
    pptx_ok = os.path.exists(active_files.get("pptx", ""))
    return jsonify({
        "excel_name": active_files["excel_name"] if excel_ok else "(none — please upload)",
        "pptx_name": active_files["pptx_name"] if pptx_ok else "(none — please upload)",
        "excel_loaded": excel_ok,
        "pptx_loaded": pptx_ok,
    })


@app.route("/api/upload-excel", methods=["POST"])
def upload_excel():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "No file selected"}), 400

    ext = os.path.splitext(f.filename)[1].lower()
    if ext not in (".xlsx", ".xlsm", ".xls"):
        return jsonify({"error": "Invalid file type. Please upload an Excel file (.xlsx, .xlsm)"}), 400

    filename = secure_filename(f.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    f.save(filepath)

    # Validate the Excel file has the expected sheets
    try:
        xls = pd.ExcelFile(filepath)
        missing = [s for s in SHEET_MAP.values() if s not in xls.sheet_names]
        if missing:
            warning = f"Warning: Missing sheets: {', '.join(missing)}. Data for these will be empty."
        else:
            warning = None
    except Exception as e:
        os.remove(filepath)
        return jsonify({"error": f"Cannot read Excel file: {str(e)}"}), 400

    active_files["excel"] = filepath
    active_files["excel_name"] = f.filename
    _save_active_files()

    return jsonify({
        "success": True,
        "filename": f.filename,
        "warning": warning,
    })


@app.route("/api/upload-pptx", methods=["POST"])
def upload_pptx():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "No file selected"}), 400

    ext = os.path.splitext(f.filename)[1].lower()
    if ext != ".pptx":
        return jsonify({"error": "Invalid file type. Please upload a .pptx file"}), 400

    filename = secure_filename(f.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    f.save(filepath)

    # Validate the PPTX file
    try:
        prs = Presentation(filepath)
        slide_count = len(prs.slides)
    except Exception as e:
        os.remove(filepath)
        return jsonify({"error": f"Cannot read PPTX file: {str(e)}"}), 400

    active_files["pptx"] = filepath
    active_files["pptx_name"] = f.filename
    _save_active_files()

    return jsonify({
        "success": True,
        "filename": f.filename,
        "slide_count": slide_count,
    })


@app.route("/api/pptx-slides")
def pptx_slides():
    """Return list of slide titles from the active PPTX template."""
    if not os.path.exists(active_files.get("pptx", "")):
        return jsonify({
            "slides": [],
            "total": 0,
            "error": "No PPT template uploaded yet. Please upload one first.",
        })
    try:
        prs = Presentation(active_files["pptx"])
        slides = []
        for i, slide in enumerate(prs.slides):
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip().split("\n")[0][:80]
                    if t:
                        title = t
                        break
            slides.append({"num": i + 1, "title": title or f"Slide {i + 1}"})
        return jsonify({"slides": slides, "total": len(slides)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/data")
def api_data():
    if not os.path.exists(active_files.get("excel", "")):
        return jsonify({
            "error": "No Excel uploaded yet. Please upload the monthly governance workbook first.",
        }), 400
    data = read_excel_data()
    summaries = compute_summaries(data)
    return jsonify({"data": data, "summaries": summaries})


@app.route("/api/generate-ppt", methods=["POST"])
def generate_ppt():
    # Pre-flight: both files must exist
    missing = []
    if not os.path.exists(active_files.get("excel", "")):
        missing.append("Excel workbook")
    if not os.path.exists(active_files.get("pptx", "")):
        missing.append("PPT template")
    if missing:
        return jsonify({
            "error": f"Please upload the {' and '.join(missing)} before generating.",
        }), 400

    req_data = request.get_json() or {}
    selected_slides = req_data.get("slides", None)

    prs = Presentation(active_files["pptx"])
    total_slides = len(prs.slides)

    if selected_slides is None:
        selected_slides = list(range(1, total_slides + 1))

    # Update automatable slides with Excel data
    try:
        data = read_excel_data()
        update_presentation(prs, data)
    except Exception as e:
        print(f"Warning: Slide data update failed: {e}")

    slides_to_remove = []
    for i in range(total_slides):
        if (i + 1) not in selected_slides:
            slides_to_remove.append(i)

    for idx in reversed(slides_to_remove):
        rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=f"LPL_QE_Connect_Generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
    )


if __name__ == "__main__":
    app.run(debug=True, port=5000)
