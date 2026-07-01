"""
Slide data updaters: populate PPT charts/tables from Excel data.

Automatable slides:
  11 - QE Highlights Execution (DailyExecution)
  12 - Automation Snapshot (InSprintData)
  13 - Defect View 1/2 (DefectData)
  14 - Defect View 2/2 (DefectData)
  16 - RAD Snapshot 1/2 (RADEnabled)
  17 - RAD Snapshot 2/2 (RADEnabled)
  18 - Primary Release Automation 1/2 (ReleaseDayTestCaseSheet)
  19 - Primary Release Automation 2/2 (ReleaseDayTestCaseSheet)
  20 - Regression Automation Trend 1/2 (TCsDetailsSheet)
  21 - Regression Automation Trend 2/2 (TCsDetailsSheet)
"""
from pptx.chart.data import CategoryChartData


def _num(v):
    try:
        return float(v) if v != "" else 0
    except (ValueError, TypeError):
        return 0


def update_chart_data(chart, categories, series_dict):
    """Update a chart with new category data. series_dict = {name: [values]}."""
    cd = CategoryChartData()
    cd.categories = categories
    for name, values in series_dict.items():
        cd.add_series(name, values)
    chart.replace_data(cd)


def _iter_shapes(shapes):
    """Recursively iterate all shapes, descending into groups."""
    for s in shapes:
        yield s
        # Group shapes have shape_type == 6 (MSO_SHAPE_TYPE.GROUP) and a .shapes attribute
        if getattr(s, "shape_type", None) == 6 and hasattr(s, "shapes"):
            yield from _iter_shapes(s.shapes)


def _get_charts_by_name(slide):
    """Return dict of shape.name -> shape for chart shapes."""
    return {s.name: s for s in slide.shapes if s.has_chart}


def _get_text_shapes(slide):
    """Return list of (shape_name, text, shape) for text shapes."""
    result = []
    for s in slide.shapes:
        if s.has_text_frame:
            result.append((s.name, s.text_frame.text.strip(), s))
    return result


def _set_text_preserve_format(shape, new_text):
    """Set text of a shape's first paragraph, preserving run formatting."""
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
    elif tf.paragraphs:
        tf.paragraphs[0].text = new_text


# ==============================================================================
# SLIDE 10: Cognizant QE Highlights (1/2) — # Scripts Distribution pie chart
# Source: MonthlySheet
#   - Filter rows where DashBoard Count > 0
#   - Sum Feature Branch Count, Develop Branch Count, Main Branch Count
# ==============================================================================
def update_slide_10(slide, data):
    """Update the #Scripts Distribution pie chart (Feature/Develop/Main) on slide 10."""
    monthly = data.get("monthly_sheet", [])
    if not monthly:
        return

    feature_sum = 0
    develop_sum = 0
    main_sum = 0
    python_total = 0
    for r in monthly:
        dash = _num(r.get("DashBoard Count", 0))
        f = _num(r.get("Feature Branch Count", 0))
        d = _num(r.get("Develop Branch Count", 0))
        m = _num(r.get("Main Branch Count", 0))
        tech = str(r.get("Technology Used", "")).strip().lower()
        # Pie chart + UI/API count: DashBoard Count > 0 AND Technology Used == Java
        if dash > 0 and tech == "java":
            feature_sum += f
            develop_sum += d
            main_sum += m
        # Python Scripts Created = sum of Feature+Develop+Main where Technology Used == Python
        if tech == "python":
            python_total += f + d + m

    total_scripts = feature_sum + develop_sum + main_sum

    for s in slide.shapes:
        if not s.has_chart:
            continue
        chart = s.chart
        try:
            cats = [str(c).strip() for c in chart.plots[0].categories]
        except Exception:
            continue
        # Match the Feature/Develop/Main pie chart by category labels
        cats_lower = [c.lower() for c in cats]
        if ("feature" in cats_lower and "develop" in cats_lower and "main" in cats_lower):
            # Preserve the existing series name from the template
            try:
                series_name = chart.plots[0].series[0].name or "Sales"
            except Exception:
                series_name = "Sales"
            # Align values to the chart's category order
            value_map = {"feature": feature_sum, "develop": develop_sum, "main": main_sum}
            values = [value_map.get(c, 0) for c in cats_lower]
            update_chart_data(chart, cats, {series_name: values})

    # Helper: find the numeric text box closest to a label text, and set its value.
    def _update_number_next_to_label(label_text, new_value):
        label_shape = None
        for s in slide.shapes:
            if s.has_text_frame and label_text in s.text_frame.text:
                label_shape = s
                break
        if label_shape is None:
            return None
        label_top = label_shape.top or 0
        label_left = label_shape.left or 0
        best = None
        best_dist = None
        for s in slide.shapes:
            if not s.has_text_frame or s is label_shape:
                continue
            txt = s.text_frame.text.strip()
            if not txt.isdigit():
                continue
            dx = (s.left or 0) - label_left
            dy = (s.top or 0) - label_top
            dist = dx * dx + dy * dy
            if best_dist is None or dist < best_dist:
                best_dist = dist
                best = s
        if best is not None:
            _set_text_preserve_format(best, str(int(new_value)))
        return best

    # UI/API Scripts Created = Feature + Develop + Main totals from the pie chart
    _update_number_next_to_label("UI/API Scripts Created", total_scripts)
    # Python Scripts Created = sum of F+D+M where Technology Used == Python
    _update_number_next_to_label("Python Scripts Created", python_total)


# ==============================================================================
# SLIDE 11: QE Highlights - Execution (DailyExecution)
# ==============================================================================
def update_slide_11(slide, data):
    """Update daily execution charts and summary numbers."""
    daily = data.get("daily_execution", [])
    if not daily:
        return

    # P0-P4 totals across ALL rows (kept available for any chart that needs it)
    p_totals_all = {"P0": 0, "P1": 0, "P2": 0, "P3": 0, "P4": 0}
    # P0-P4 totals filtered by Daily/Weekly — used for Critical/Core/All and P0-P4 charts
    p_totals = {"P0": 0, "P1": 0, "P2": 0, "P3": 0, "P4": 0}
    total_passed = 0
    total_failed = 0
    # "Testcases configured in nightly batch" = sum of P0+P1+P2 where
    # Automation Execution Frequency is Daily or Weekly
    nightly_batch_total = 0
    # Pass rate = sum(P0-P2 Passed TCs) / sum(P0+P1+P2) across rows
    # where Automation Execution Frequency is Daily or Weekly
    p0_p2_total = 0
    p0_p2_passed = 0
    for r in daily:
        p0 = _num(r.get("P0", 0))
        p1 = _num(r.get("P1", 0))
        p2 = _num(r.get("P2", 0))
        p3 = _num(r.get("P3", 0))
        p4 = _num(r.get("P4", 0))
        p_totals_all["P0"] += p0
        p_totals_all["P1"] += p1
        p_totals_all["P2"] += p2
        p_totals_all["P3"] += p3
        p_totals_all["P4"] += p4
        total_passed += _num(r.get("Total Passed TC #", 0))
        total_failed += _num(r.get("Total Failed TC #", 0))
        freq = str(r.get("Automation Execution Frequency ",
                          r.get("Automation Execution Frequency", ""))).strip().lower()
        if freq in ("daily", "weekly"):
            p_totals["P0"] += p0
            p_totals["P1"] += p1
            p_totals["P2"] += p2
            p_totals["P3"] += p3
            p_totals["P4"] += p4
            nightly_batch_total += p0 + p1 + p2
            p0_p2_total += p0 + p1 + p2
            p0_p2_passed += _num(r.get("P0 - P2 Passed TCs", 0))

    # Critical/Core/All on the bar chart, filtered by Daily/Weekly
    total_configured = sum(p_totals.values())
    critical = p_totals["P0"] + p_totals["P1"]
    core = critical + p_totals["P2"]

    charts = _get_charts_by_name(slide)

    # Chart 26: Critical/Core/All bar chart
    for name, shape in charts.items():
        chart = shape.chart
        cats = [str(c).strip() for c in chart.plots[0].categories]
        if "Critical" in cats[0] if cats else False:
            update_chart_data(chart, ["Critical", "Core", "All"],
                              {"Daily Execution": [critical, core, total_configured]})
        elif "P0" in cats[0] if cats else False:
            update_chart_data(chart, ["P0", "P1", "P2", "P3", "P4"],
                              {"Daily Execution": [p_totals["P0"], p_totals["P1"], p_totals["P2"],
                                                   p_totals["P3"], p_totals["P4"]]})

    # Update "Testcases configured in nightly batch" number by finding the
    # integer text box closest to that label (recursively descends into groups).
    all_shapes = list(_iter_shapes(slide.shapes))
    label_shape = None
    for s in all_shapes:
        if s.has_text_frame and "nightly batch" in s.text_frame.text.lower():
            label_shape = s
            break
    if label_shape is not None:
        label_top = label_shape.top or 0
        label_left = label_shape.left or 0
        best = None
        best_dist = None
        for s in all_shapes:
            if not s.has_text_frame or s is label_shape:
                continue
            txt = s.text_frame.text.strip()
            if not txt.isdigit():
                continue
            dx = (s.left or 0) - label_left
            dy = (s.top or 0) - label_top
            dist = dx * dx + dy * dy
            if best_dist is None or dist < best_dist:
                best_dist = dist
                best = s
        if best is not None:
            _set_text_preserve_format(best, str(int(nightly_batch_total)))

    # Pass rate = sum(P0-P2 Passed TCs) / sum(P0+P1+P2) * 100 (Daily/Weekly only)
    pass_rate = round(p0_p2_passed / p0_p2_total * 100) if p0_p2_total else 0
    # Find the "~NN%" text box (near "Pass rate" label) — descend into groups.
    pass_label = None
    for s in all_shapes:
        if s.has_text_frame and "Pass rate" in s.text_frame.text:
            pass_label = s
            break
    if pass_label is not None:
        label_top = pass_label.top or 0
        label_left = pass_label.left or 0
        best = None
        best_dist = None
        for s in all_shapes:
            if not s.has_text_frame or s is pass_label:
                continue
            txt = s.text_frame.text.strip()
            if "%" not in txt:
                continue
            dx = (s.left or 0) - label_left
            dy = (s.top or 0) - label_top
            dist = dx * dx + dy * dy
            if best_dist is None or dist < best_dist:
                best_dist = dist
                best = s
        if best is not None:
            _set_text_preserve_format(best, f"~{pass_rate}%")


# ==============================================================================
# SLIDE 12: Automation Snapshot (InSprintData)
# ==============================================================================
def _automation_coverage_pct(v):
    """Normalize an Automation coverage cell to a 0..100 percentage.

    Accepts 0.85, '85%', '85', 85, etc. Returns None if unparseable.
    """
    if v is None or v == "":
        return None
    try:
        f = float(v)
    except (TypeError, ValueError):
        s = str(v).strip().replace("%", "")
        try:
            f = float(s)
        except ValueError:
            return None
    # Excel often stores percentages as fractions (0.85 = 85%)
    return f * 100 if abs(f) <= 1 else f


def _set_cell_text_preserve(cell, new_text):
    """Replace a table cell's entire contents with `new_text` while keeping
    the first run's font/color/size formatting. Any additional paragraphs
    left over from the template are removed so old text doesn't leak through.
    """
    from pptx.oxml.ns import qn
    tf = cell.text_frame
    txBody = tf._txBody
    # Remove every paragraph after the first
    paragraphs = txBody.findall(qn("a:p"))
    for p in paragraphs[1:]:
        txBody.remove(p)
    # Set the first paragraph's text, preserving its first run's formatting
    first_p = tf.paragraphs[0]
    if first_p.runs:
        first_p.runs[0].text = new_text
        for run in first_p.runs[1:]:
            run.text = ""
    else:
        first_p.text = new_text


def _classify_insprint_teams(data):
    """Classify each Scrum team into one of 5 automation-health buckets based
    on rows from the LAST 2 sprint cycles in InSprintData.

    Aggregates per team across those rows:
      max_cov  = max(Automation coverage)  (percentage, 0..100)
      sum_td   = sum(Auto tech debt P0..P4)
      sum_mnt  = sum(TCs Maintained)
      comments = list of non-blank 'Comment (If Automation Cov is 0%)' values

    Priority (first matching wins):
      max_cov >= 70                                    -> 'gt70'   (>70 % In-Sprint)
      0 < max_cov < 70                                 -> 'lt70'   (<70 % In-Sprint)
      max_cov == 0 AND sum_td > 0                      -> 'tech'   (Tech-Debt Automation)
      max_cov == 0 AND sum_td == 0 AND sum_mnt > 0     -> 'maint'  (Script Maintainance)
      max_cov == 0 AND sum_td == 0 AND sum_mnt == 0    -> 'none'   (No Automation)

    Returns (buckets, team_agg, sprints)
      buckets: {'gt70': [teams], 'lt70': [...], 'tech': [...], 'maint': [...], 'none': [...]}
      team_agg: {team: {'max_cov', 'sum_td', 'sum_mnt', 'comments'}}
      sprints:  list of the 2 sprint names actually used (for reporting)
    """
    insprint = data.get("insprint_data", [])
    if not insprint:
        return {"gt70": [], "lt70": [], "tech": [], "maint": [], "none": []}, {}, []

    # Detect Sprint column (may have trailing whitespace variants)
    sprint_col = "Sprint"
    for k in insprint[0].keys():
        if str(k).strip().lower() == "sprint":
            sprint_col = k
            break
    sprints = sorted(
        {str(r.get(sprint_col, "")).strip() for r in insprint
         if str(r.get(sprint_col, "")).strip().startswith("Sprint ")},
        key=_sprint_sort_key,
    )[-2:]
    if not sprints:
        return {"gt70": [], "lt70": [], "tech": [], "maint": [], "none": []}, {}, []

    team_agg = {}
    for r in insprint:
        if str(r.get(sprint_col, "")).strip() not in sprints:
            continue
        team = str(r.get("Scrum team", "")).strip()
        if not team:
            continue
        cov = _automation_coverage_pct(r.get("Automation coverage"))
        td = sum(_num(r.get(f"Auto tech debt P{i}", 0)) for i in range(5))
        mnt = _num(r.get("TCs Maintained", 0))
        comment = str(r.get("Comment (If Automation Cov is 0%)", "") or "").strip()
        d = team_agg.setdefault(team, {"max_cov": None, "sum_td": 0.0,
                                       "sum_mnt": 0.0, "comments": []})
        if cov is not None and (d["max_cov"] is None or cov > d["max_cov"]):
            d["max_cov"] = cov
        d["sum_td"] += td
        d["sum_mnt"] += mnt
        if comment and comment not in d["comments"]:
            d["comments"].append(comment)

    buckets = {"gt70": [], "lt70": [], "tech": [], "maint": [], "none": []}
    for team, d in team_agg.items():
        mc = d["max_cov"] if d["max_cov"] is not None else 0
        if mc >= 70:
            buckets["gt70"].append(team)
        elif mc > 0:
            buckets["lt70"].append(team)
        elif d["sum_td"] > 0:
            buckets["tech"].append(team)
        elif d["sum_mnt"] > 0:
            buckets["maint"].append(team)
        else:
            buckets["none"].append(team)
    return buckets, team_agg, sprints


def update_slide_12(slide, data):
    """Update the Automation Snapshot bar chart AND the 'Reason for no
    Automation' table on slide 12.

    Chart (5 bars, restricted to last 2 sprint cycles, per-team classification):
      '>70 % In-Sprint'      = teams whose max Automation coverage >= 70
      '<70 % In-Sprint'      = teams whose max Automation coverage is > 0 and < 70
      'Tech-Debt Automation' = teams with Automation coverage = 0 and
                               sum(Auto tech debt P0..P4) > 0
      'Script Maintainance'  = teams with Automation coverage = 0 and
                               sum(Auto tech debt P0..P4) = 0 and TCs Maintained > 0
      'No Automation'        = teams with Automation coverage = 0 and
                               sum(Auto tech debt P0..P4) = 0 and TCs Maintained = 0

    Table 'Reason for no Automation':
      Only teams in the 'No Automation' bucket. Grouped by their
      'Comment (If Automation Cov is 0%)' value. One row per distinct
      reason, showing '<reason> | [team1, team2, ...]' and the team count.
    """
    buckets, team_agg, sprints = _classify_insprint_teams(data)
    if not any(buckets.values()):
        return

    # ---- Update the slide title with the actual sprint range ----
    # Template ships something like "Automation Snapshot - Sprint 25.2.3 - Sprint 25.2.4";
    # rewrite it using the sprints we actually classified.
    if len(sprints) >= 1:
        sprint_range = (sprints[0] if len(sprints) == 1
                        else f"{sprints[0]} - {sprints[1]}")
        for s in slide.shapes:
            if s.has_text_frame and "Automation Snapshot" in s.text_frame.text:
                _set_text_preserve_format(s, f"Automation Snapshot - {sprint_range}")
                break

    # ---- Update the 5-bar chart ----
    counts_by_cat = {
        "No Automation":        len(buckets["none"]),
        "Script Maintainance":  len(buckets["maint"]),
        "Tech-Debt Automation": len(buckets["tech"]),
        "<70 % In-Sprint":      len(buckets["lt70"]),
        ">70 % In-Sprint":      len(buckets["gt70"]),
    }
    for s in slide.shapes:
        if not s.has_chart:
            continue
        chart = s.chart
        try:
            cats = [str(c).strip() for c in chart.plots[0].categories]
        except Exception:
            continue
        if "<70 % In-Sprint" not in cats or ">70 % In-Sprint" not in cats:
            continue
        try:
            series_name = chart.plots[0].series[0].name or "Teams"
        except Exception:
            series_name = "Teams"
        new_values = [counts_by_cat.get(c, 0) for c in cats]
        update_chart_data(chart, cats, {series_name: new_values})
        break

    # ---- Update the 'Reason for no Automation' table ----
    # Group No-Automation teams by their reason comment
    reason_to_teams = {}
    for team in buckets["none"]:
        d = team_agg.get(team, {"comments": []})
        reason = d["comments"][0] if d["comments"] else "(No reason given)"
        reason_to_teams.setdefault(reason, []).append(team)

    for s in slide.shapes:
        if not getattr(s, "has_table", False):
            continue
        tbl = s.table
        # Expect: header row + N data rows + Total row.
        n_rows = len(tbl.rows)
        if n_rows < 3:
            continue
        data_rows_capacity = n_rows - 2  # exclude header + total
        # Sort reasons by team count desc, then alpha
        entries = sorted(reason_to_teams.items(),
                         key=lambda x: (-len(x[1]), x[0].lower()))
        # If more reasons than capacity, lump the tail into 'Others'
        if len(entries) > data_rows_capacity:
            head = entries[:data_rows_capacity - 1]
            tail = entries[data_rows_capacity - 1:]
            others_teams = [t for _, teams in tail for t in teams]
            others_reasons = "; ".join(r for r, _ in tail)
            head.append((f"Others ({others_reasons})", others_teams))
            entries = head
        grand = sum(len(t) for _, t in entries)
        for i in range(data_rows_capacity):
            row_idx = i + 1
            if i < len(entries):
                reason, teams = entries[i]
                cell_text = f"{reason} | [{', '.join(teams)}]"
                _set_cell_text_preserve(tbl.cell(row_idx, 0), cell_text)
                _set_cell_text_preserve(tbl.cell(row_idx, 1), str(len(teams)))
            else:
                _set_cell_text_preserve(tbl.cell(row_idx, 0), "")
                _set_cell_text_preserve(tbl.cell(row_idx, 1), "")
        # Total row (last row)
        total_idx = n_rows - 1
        _set_cell_text_preserve(tbl.cell(total_idx, 0), "Total")
        _set_cell_text_preserve(tbl.cell(total_idx, 1), str(grand))
        break


# ==============================================================================
# SLIDES 13-14: Defect View (DefectData)
# ==============================================================================

# Chart-to-domain mapping based on PPT structure
DEFECT_SLIDE_13_DOMAINS = {
    "Chart 5": "ALM-Technology",
    "Chart 7": "Corporate Systems (TFG)",
    "Chart 9": "Custody, Clearing & Settlement",
    "Chart 11": "Data/Platform Modernization",
    "Chart 13": "Infosec",
    "Chart 15": "Practice Management",
}
DEFECT_SLIDE_14_DOMAINS = {
    "Chart 5": "Service and Support",
    "Chart 7": "SRC",
    "Chart 9": "Technology",
    "Chart 11": "Trading",
}
# Label shape -> domain mapping for the text labels showing "Domain (count)"
DEFECT_LABEL_13 = {
    "Rectangle 4": "ALM-Technology",
    "Rectangle 6": "Corporate Systems (TFG)",
    "Rectangle 8": "Custody, Clearing & Settlement",
    "Rectangle 10": "Data/Platform Modernization",
    "Rectangle 12": "Infosec",
    "Rectangle 14": "Practice Management",
}
DEFECT_LABEL_14 = {
    "Rectangle 4": "Service and Support",
    "Rectangle 6": "SRC",
    "Rectangle 8": "Technology",
    "Rectangle 10": "Trading",
    "Rectangle 17": "Automation",
}

# Per-slide config for the Manual defect mini-charts (slides 13 & 14):
#   (label_shape, chart_shape, excel_sub_domain, display_label)
DEFECT_VIEW_SLIDE_13 = [
    ("Rectangle 4",  "Chart 5",  "ALM-Technology",                 "ALM – Technology"),
    ("Rectangle 6",  "Chart 7",  "Corporate Systems(TFG)",         "Corporate Systems (TFG)"),
    ("Rectangle 8",  "Chart 9",  "Custody, Clearing & Settlement", "Custody, Clearing & Settlement"),
    ("Rectangle 10", "Chart 11", "Data",                           "Data"),
    ("Rectangle 12", "Chart 13", "Infosec",                        "Infosec"),
    ("Rectangle 14", "Chart 15", "Practice Management",            "Practice Management"),
]
DEFECT_VIEW_SLIDE_14 = [
    ("Rectangle 4",  "Chart 5",  "Service and Support", "Service and Support"),
    ("Rectangle 6",  "Chart 7",  "SRC",                 "SRC"),
    ("Rectangle 8",  "Chart 9",  "Technology",          "Technology"),
    ("Rectangle 10", "Chart 11", "Trading",             "Trading"),
]


def _sprint_sort_key(sprint_name):
    """Numeric sort key for sprint names like 'Sprint 26.1.4'.

    Extracts every integer in the string and returns them as a tuple so
    'Sprint 26.1.10' correctly sorts AFTER 'Sprint 26.1.2' (and any future
    naming variant — 'Sprint 27.2.1', 'Sprint 2026.1.4', etc. — keeps working).
    Falls back to the raw string for anything without numbers.
    """
    import re
    nums = re.findall(r"\d+", str(sprint_name))
    return tuple(int(n) for n in nums) if nums else (str(sprint_name),)


def _last_two_sprint_cycles(data):
    """Return the two most-recent sprint-cycle values from DefectData in which
    *at least one row has a non-empty Sub Domain* (excluding 'TBD').

    A sprint cycle is any Sprint value starting with 'Sprint ' (e.g. 'Sprint 26.1.4').
    Monthly buckets like \"January'26\" are excluded. Sprints are sorted
    numerically (so 'Sprint 26.1.10' comes after 'Sprint 26.1.2', not before).

    Sprints whose rows all have a blank or TBD Sub Domain are skipped —
    otherwise the Defect View slides would silently produce empty charts
    just because the latest sprint hasn't had Sub Domain filled in yet.
    """
    sprints_with_subdomain = set()
    for r in data.get("defect_data", []):
        s = str(r.get("Sprint", "")).strip()
        if not s.startswith("Sprint "):
            continue
        sub = str(r.get("Sub Domain", "")).strip()
        if sub and sub != "TBD":
            sprints_with_subdomain.add(s)
    return sorted(sprints_with_subdomain, key=_sprint_sort_key)[-2:]


def _compute_automation_prod_defects(data):
    """{sub_domain: total_defects} for the Slide 14 Automation panel.

    Filters:
      Sprint in the LAST 2 sprint cycles that have Sub Domain data
        (same window used by the per-Sub Domain Manual charts on Slide 13/14),
      Type Of Defect == Automation,
      Sub Domain != TBD (and not blank).

    Per Sub Domain, combine every severity across both InSprint and Regression:
      total = InSprint (Fatel + Serious + Medium + Low)
            + Regression (Fatel + Serious + Medium + Low)
    """
    sprints = set(_last_two_sprint_cycles(data))
    if not sprints:
        return {}
    severity_cols = [
        "InSprint Fatel",   "Regression Fatel",
        "InSprint Serious", "Regression Serious",
        "InSprint Medium",  "Regression Medium",
        "InSprint Low",     "Regression Low",
    ]
    result = {}
    for r in data.get("defect_data", []):
        s = str(r.get("Sprint", "")).strip()
        if s not in sprints:
            continue
        if str(r.get("Type Of Defect", "")).strip().lower() != "automation":
            continue
        sub = str(r.get("Sub Domain", "")).strip()
        if not sub or sub == "TBD":
            continue
        result[sub] = result.get(sub, 0) + sum(_num(r.get(c, 0)) for c in severity_cols)
    return result


# Maps the labels shown in the Automation chart (slide 14, Chart 18) to the
# Excel Sub Domain values. Keep the order in sync with the chart's existing
# categories so legend colors stay stable.
AUTOMATION_CHART_LABELS = [
    ("SRC",                             "SRC"),
    ("Service and Support",             "Service and Support"),
    ("Trading",                         "Trading"),
    ("Corporate Systems",               "Corporate Systems(TFG)"),
    ("Custody, clearing & settlement",  "Custody, Clearing & Settlement"),
    ("Data",                            "Data"),
    ("Infosec",                         "Infosec"),
    ("Practice Management",             "Practice Management"),
    ("Technology",                      "Technology"),
    ("ALM",                             "ALM-Technology"),
]


def _compute_manual_defects_by_subdomain(data, sprints):
    """{sub_domain: {Fatal, Serious, Medium, Low}} filtered by:
       Sprint ∈ sprints, Type Of Defect == Manual, Sub Domain != TBD.
    """
    result = {}
    sprint_set = set(sprints)
    for r in data.get("defect_data", []):
        if str(r.get("Sprint", "")).strip() not in sprint_set:
            continue
        if str(r.get("Type Of Defect", "")).strip().lower() != "manual":
            continue
        sub = str(r.get("Sub Domain", "")).strip()
        if not sub or sub == "TBD":
            continue
        if sub not in result:
            result[sub] = {"Fatal": 0, "Serious": 0, "Medium": 0, "Low": 0}
        result[sub]["Fatal"]   += _num(r.get("InSprint Fatel", 0))   + _num(r.get("Regression Fatel", 0))
        result[sub]["Serious"] += _num(r.get("InSprint Serious", 0)) + _num(r.get("Regression Serious", 0))
        result[sub]["Medium"]  += _num(r.get("InSprint Medium", 0))  + _num(r.get("Regression Medium", 0))
        result[sub]["Low"]     += _num(r.get("InSprint Low", 0))     + _num(r.get("Regression Low", 0))
    return result


def _compute_defect_by_domain(data, defect_type=None):
    """Compute {domain: {Fatal, Serious, Medium, Low}} from DefectData."""
    result = {}
    for r in data.get("defect_data", []):
        if defect_type and str(r.get("Type Of Defect", "")) != defect_type:
            continue
        domain = str(r.get("Domain", "Unknown"))
        if domain not in result:
            result[domain] = {"Fatal": 0, "Serious": 0, "Medium": 0, "Low": 0}
        result[domain]["Fatal"] += _num(r.get("InSprint Fatel", 0)) + _num(r.get("Regression Fatel", 0))
        result[domain]["Serious"] += _num(r.get("InSprint Serious", 0)) + _num(r.get("Regression Serious", 0))
        result[domain]["Medium"] += _num(r.get("InSprint Medium", 0)) + _num(r.get("Regression Medium", 0))
        result[domain]["Low"] += _num(r.get("InSprint Low", 0)) + _num(r.get("Regression Low", 0))
    return result


def _domain_match(excel_domain, ppt_domain):
    """Fuzzy match domain names between Excel and PPT."""
    ed = excel_domain.lower().strip()
    pd = ppt_domain.lower().strip()
    if ed == pd:
        return True
    # Common mappings
    mappings = {
        "alm-technology": ["alm", "alm-technology", "account lifecycle management"],
        "corporate systems (tfg)": ["corporate systems", "corporate systems(tfg)", "corporate systems (tfg)"],
        "custody, clearing & settlement": ["custody", "custody, cleaning & settlement", "custody, clearing & settlement"],
        "data/platform modernization": ["data", "data/platform modernization"],
        "service and support": ["service and support", "service & support"],
        "advisor experience": ["advisor experience", "src"],
        "practice management": ["practice management"],
        "investor experience": ["investor experience"],
        "trading": ["trading"],
        "infosec": ["infosec"],
        "technology": ["technology"],
    }
    for key, variants in mappings.items():
        if pd in variants and ed in variants:
            return True
        if pd == key and ed in variants:
            return True
        if ed == key and pd in variants:
            return True
    return False


def _find_domain_data(domain_defects, ppt_domain):
    """Find matching domain data using fuzzy matching."""
    for ed, vals in domain_defects.items():
        if _domain_match(ed, ppt_domain):
            return vals
    return {"Fatal": 0, "Serious": 0, "Medium": 0, "Low": 0}


def update_defect_slide(slide, data, chart_domain_map, label_domain_map):
    """Update a defect view slide's charts and labels."""
    domain_defects = _compute_defect_by_domain(data)
    auto_defects = _compute_defect_by_domain(data, defect_type="Automation")

    charts = _get_charts_by_name(slide)
    for chart_name, ppt_domain in chart_domain_map.items():
        if chart_name not in charts:
            continue
        chart = charts[chart_name].chart
        vals = _find_domain_data(domain_defects, ppt_domain)
        update_chart_data(chart, ["Fatal", "Serious", "Medium", "Low"],
                          {"Defect Metrics": [vals["Fatal"], vals["Serious"], vals["Medium"], vals["Low"]]})

    # Handle automation defect chart on slide 14 (Chart 18)
    if "Chart 18" in charts:
        chart = charts["Chart 18"].chart
        auto_domains = ["SRC", "Service and Support", "Trading", "Corporate Systems (TFG)",
                        "Custody, Clearing & Settlement", "Data/Platform Modernization",
                        "Infosec", "Practice Management", "Technology", "ALM-Technology"]
        auto_vals = []
        for d in auto_domains:
            dv = _find_domain_data(auto_defects, d)
            total = dv["Fatal"] + dv["Serious"] + dv["Medium"] + dv["Low"]
            auto_vals.append(total if total else None)
        display_names = ["SRC", "Service and Support", "Trading", "Corporate Systems",
                         "Custody, clearing & settlement", "Data", "Infosec",
                         "Practice Management", "Technology", "ALM"]
        update_chart_data(chart, display_names, {"": auto_vals})

    # Update domain labels with counts
    for shape_name, text, shape in _get_text_shapes(slide):
        if shape_name in label_domain_map:
            ppt_domain = label_domain_map[shape_name]
            if ppt_domain == "Automation":
                total = sum(
                    sum(v.values()) for v in auto_defects.values()
                )
                _set_text_preserve_format(shape, f"Automation ({int(total)})")
            else:
                vals = _find_domain_data(domain_defects, ppt_domain)
                total = int(vals["Fatal"] + vals["Serious"] + vals["Medium"] + vals["Low"])
                # Extract display name from PPT domain
                display = ppt_domain.split("/")[0] if "/" in ppt_domain else ppt_domain
                # Shorten some names
                display = display.replace("Account Lifecycle Management", "ALM-Technology")
                _set_text_preserve_format(shape, f"{display} ({total})")


def _update_defect_view_slide(slide, data, slide_config):
    """Populate a Defect View slide (13 or 14).

    For each (label, chart, sub_domain, display) row in slide_config:
      - update the chart with [Fatal, Serious, Medium, Low] for that sub-domain
        (filtered to the last 2 sprint cycles, Manual defects only, no TBD)
      - update the label with '<display> (<total>)'

    Also updates the slide title with the actual sprint range, e.g.
    'Defect VIEW – Sprint 26.1.3 - Sprint 26.1.4'.
    """
    sprints = _last_two_sprint_cycles(data)
    if not sprints:
        return
    defects = _compute_manual_defects_by_subdomain(data, sprints)

    shape_map = {s.name: s for s in slide.shapes}

    for label_name, chart_name, sub_domain, display in slide_config:
        vals = defects.get(sub_domain, {"Fatal": 0, "Serious": 0, "Medium": 0, "Low": 0})
        total = int(vals["Fatal"] + vals["Serious"] + vals["Medium"] + vals["Low"])

        # Update label "<display> (<total>)"
        if label_name in shape_map:
            _set_text_preserve_format(shape_map[label_name], f"{display} ({total})")

        # Update chart values
        if chart_name in shape_map and shape_map[chart_name].has_chart:
            chart = shape_map[chart_name].chart
            try:
                series_name = chart.plots[0].series[0].name or "Defect Metrics"
            except Exception:
                series_name = "Defect Metrics"
            update_chart_data(chart, ["Fatal", "Serious", "Medium", "Low"], {
                series_name: [
                    int(vals["Fatal"]),
                    int(vals["Serious"]),
                    int(vals["Medium"]),
                    int(vals["Low"]),
                ],
            })

    # Update the slide title with the dynamic sprint range
    sprint_range = (sprints[0] if len(sprints) == 1
                    else f"{sprints[0]} - {sprints[1]}")
    for s in slide.shapes:
        if s.has_text_frame and "Defect VIEW" in s.text_frame.text:
            _set_text_preserve_format(s, f"Defect VIEW – {sprint_range}")
            break


def update_slide_13(slide, data):
    _update_defect_view_slide(slide, data, DEFECT_VIEW_SLIDE_13)
    # Keep prior automation/extra logic out of slide 13 — only the per-sub-domain
    # Manual defect mini-charts are populated, as requested.


def update_slide_14(slide, data):
    _update_defect_view_slide(slide, data, DEFECT_VIEW_SLIDE_14)

    # Automation panel on slide 14 (Rectangle 17 + Chart 18):
    # Sum 'Prod Defect Count' across rows where
    #   Sprint is a monthly bucket (e.g. "January'26")  AND
    #   Type Of Defect == Automation                    AND
    #   Sub Domain != TBD
    # grouped by Sub Domain. Chart categories are kept in the template's
    # order; missing sub-domains are written as None (so empty bars hide).
    auto_by_sub = _compute_automation_prod_defects(data)
    shape_map = {s.name: s for s in slide.shapes}

    if "Chart 18" in shape_map and shape_map["Chart 18"].has_chart:
        chart = shape_map["Chart 18"].chart
        try:
            series_name = chart.plots[0].series[0].name or "Defect Metrics"
        except Exception:
            series_name = "Defect Metrics"
        display_labels = [d for d, _ in AUTOMATION_CHART_LABELS]
        values = []
        for _, sub_key in AUTOMATION_CHART_LABELS:
            v = int(auto_by_sub.get(sub_key, 0))
            values.append(v if v > 0 else None)
        update_chart_data(chart, display_labels, {series_name: values})

    if "Rectangle 17" in shape_map:
        total = int(sum(auto_by_sub.values()))
        _set_text_preserve_format(shape_map["Rectangle 17"], f"Automation ({total})")


# ==============================================================================
# SLIDES 16-17: RAD Snapshot (RADEnabled)
# ==============================================================================

# Panel config for each RAD slide:
#   (label_shape, chart_shape, lookup_key, display_label, by_team)
# When by_team is False the key is a Sub Domain (Excel 'Sub Domain' column).
# When by_team is True the key is a Scrum Team Name (used only for special
# team-level panels such as 'TMLB Muppets' whose Sub Domain is blank).
RAD_SLIDE_16_PANELS = [
    ("TextBox 26", "Chart 60", "Service and Support",     "Service and Support",     False),
    ("TextBox 43", "Chart 44", "SRC",                     "SRC",                     False),
    ("TextBox 56", "Chart 57", "ALM-Technology",          "ALM-Technology",          False),
    ("TextBox 5",  "Chart 4",  "Practice Management",     "Practice Management",     False),
    ("TextBox 34", "Chart 33", "Trading",                 "Trading",                 False),
    ("TextBox 7",  "Chart 6",  "Corporate Systems(TFG)",  "Corporate Systems(TFG)",  False),
]
RAD_SLIDE_17_PANELS = [
    ("TextBox 26", "Chart 60", "Custody, Clearing & Settlement", "Custody, Clearing & Settlement", False),
    ("TextBox 43", "Chart 44", "Data",                            "Data",                           False),
    ("TextBox 56", "Chart 57", "TMLB Muppets",                    "TMLB Muppets",                   True),
    ("TextBox 2",  "Chart 8",  "Home Office",                     "Home Office",                   False),
    ("TextBox 9",  "Chart 10", "Infosec",                         "Infosec",                       False),
    ("TextBox 11", "Chart 12", "Technology",                      "Technology",                    False),
]


def _compute_rad_by_subdomain(data):
    """Group RADEnabled rows by Sub Domain (skipping TBD / blank Sub Domain).

    Returns {sub_domain: {isolated, qa_pull, rad_enabled, in_progress}} where:
      isolated     = sum(Independent Build #)
      qa_pull      = sum(# Builds QA Pull Enabled)
      rad_enabled  = sum(# Builds RAD Enabled)
      in_progress  = max(0, qa_pull - rad_enabled)
    """
    by_sub = {}
    for r in data.get("rad_enabled", []):
        sub = str(r.get("Sub Domain", "")).strip()
        if not sub or sub == "TBD":
            continue
        if sub not in by_sub:
            by_sub[sub] = {"isolated": 0, "qa_pull": 0, "rad_enabled": 0, "in_progress": 0}
        by_sub[sub]["isolated"]    += _num(r.get("Independent Build #", 0))
        by_sub[sub]["qa_pull"]     += _num(r.get("# Builds QA Pull Enabled", 0))
        by_sub[sub]["rad_enabled"] += _num(r.get("# Builds RAD Enabled", 0))
    # In-progress is QA Pull minus RAD Enabled (floored at 0)
    for v in by_sub.values():
        v["in_progress"] = max(0, v["qa_pull"] - v["rad_enabled"])
    return by_sub


def _compute_rad_by_team(data):
    """Per-team RAD aggregates — used for special team-level panels
    (e.g. 'TMLB Muppets' whose Sub Domain is blank in the source data).
    """
    by_team = {}
    for r in data.get("rad_enabled", []):
        team = str(r.get("Scrum Team Name", "")).strip()
        if not team:
            continue
        if team not in by_team:
            by_team[team] = {"isolated": 0, "qa_pull": 0, "rad_enabled": 0, "in_progress": 0}
        by_team[team]["isolated"]    += _num(r.get("Independent Build #", 0))
        by_team[team]["qa_pull"]     += _num(r.get("# Builds QA Pull Enabled", 0))
        by_team[team]["rad_enabled"] += _num(r.get("# Builds RAD Enabled", 0))
    for v in by_team.values():
        v["in_progress"] = max(0, v["qa_pull"] - v["rad_enabled"])
    return by_team


def _update_rad_slide(slide, data, panels):
    """Populate one RAD Snapshot slide given its panel config."""
    rad_sub = _compute_rad_by_subdomain(data)
    rad_team = _compute_rad_by_team(data)
    shape_map = {s.name: s for s in slide.shapes}

    cats = ["Isolated Build", "QA Pull Approved", "RAD Enable ", "RAD Enablement In-Progress"]

    for label_name, chart_name, key, display, by_team in panels:
        vals = rad_team.get(key) if by_team else rad_sub.get(key)
        vals = vals or {"isolated": 0, "qa_pull": 0, "rad_enabled": 0, "in_progress": 0}

        # Chart
        if chart_name in shape_map and shape_map[chart_name].has_chart:
            chart = shape_map[chart_name].chart
            try:
                series_name = chart.plots[0].series[0].name or "Series 1"
            except Exception:
                series_name = "Series 1"
            update_chart_data(chart, cats, {series_name: [
                int(vals["isolated"]),
                int(vals["qa_pull"]),
                int(vals["rad_enabled"]),
                int(vals["in_progress"]),
            ]})

        # Label with percentage = RAD Enabled / Independent Build
        if label_name in shape_map:
            pct = round(vals["rad_enabled"] / vals["isolated"] * 100) if vals["isolated"] else 0
            _set_text_preserve_format(shape_map[label_name], f"{display} ~ {pct} % RAD Enabled")


# ==============================================================================
# SLIDE 15: Production & Post-Production Defect View
#   Source: DefectData, monthly buckets only (Sprint != 'Sprint X.Y.Z'),
#           Type Of Defect in {Manual, Automation}, Sub Domain != TBD,
#           Prod Defect Count > 0.
#   Chart 12 = Domain Wise Defects (per-Sub Domain share of total Prod Defects)
#   Title    = "PRODUCTION & Post-production Defect View – <months>"
# ==============================================================================

# Maps the Chart 12 category labels (as they appear in the template) to the
# Excel Sub Domain values.
SLIDE_15_DOMAIN_LABELS = [
    ("ALM-Technology",                   "ALM-Technology"),
    ("Corporate Systems (TFG)",          "Corporate Systems(TFG)"),
    ("Custody, Cleaning & Settlements",  "Custody, Clearing & Settlement"),
    ("Data",                             "Data"),
    ("Infosec",                          "Infosec"),
    ("Practice Management",              "Practice Management"),
    ("Service and Support",              "Service and Support"),
    ("SRC",                              "SRC"),
    ("Technology",                       "Technology"),
    ("Trading",                          "Trading"),
]

# Short-form month names used in the dynamic title
_MONTH_SHORT = {
    "january": "Jan", "february": "Feb", "march": "Mar", "april": "Apr",
    "may": "May", "june": "Jun", "july": "Jul", "august": "Aug",
    "september": "Sep", "october": "Oct", "november": "Nov", "december": "Dec",
}


def _prod_defects_by_subdomain(data):
    """{sub_domain: prod_defect_count_total} for slide 15.

    Filters:
      Sprint is a monthly bucket (NOT 'Sprint X.Y.Z'),
      Type Of Defect in {Manual, Automation},
      Sub Domain != TBD.
    """
    result = {}
    for r in data.get("defect_data", []):
        s = str(r.get("Sprint", "")).strip()
        if not s or s.startswith("Sprint "):
            continue
        t = str(r.get("Type Of Defect", "")).strip().lower()
        if t not in ("manual", "automation"):
            continue
        sub = str(r.get("Sub Domain", "")).strip()
        if not sub or sub == "TBD":
            continue
        result[sub] = result.get(sub, 0) + _num(r.get("Prod Defect Count", 0))
    return result


def _monthly_buckets_in_data(data):
    """Return monthly Sprint values in encounter order, e.g. [\"January'26\", \"February'26\"]."""
    seen = []
    for r in data.get("defect_data", []):
        s = str(r.get("Sprint", "")).strip()
        if not s or s.startswith("Sprint "):
            continue
        if s not in seen:
            seen.append(s)
    # Sort by (year, month-index) so display order is chronological regardless of row order
    def _key(b):
        # Format like "February'26"
        try:
            name, yr = b.split("'")
            return (int(yr), list(_MONTH_SHORT.keys()).index(name.strip().lower()))
        except Exception:
            return (9999, 99)
    return sorted(seen, key=_key)


def _format_month_range(monthly_buckets):
    """Format a list of monthly bucket values into a short title, e.g. \"Jan & Feb'26\"."""
    if not monthly_buckets:
        return ""
    parts = []
    year = None
    for b in monthly_buckets:
        try:
            name, yr = b.split("'")
            short = _MONTH_SHORT.get(name.strip().lower(), name.strip()[:3])
            parts.append(short)
            year = yr
        except Exception:
            parts.append(b)
    if year:
        return f"{' & '.join(parts)}'{year}"
    return " & ".join(parts)


def _qa_miss_split(data):
    """Return (qa_miss_total, not_qa_miss_total) summed over Prod Defect Count.

    A row is QA Miss iff its 'Prod Defect RCA' text contains the literal
    substring 'qa miss' (case-insensitive); otherwise Not QA Miss.
    Same row filters as Chart 12: monthly bucket, Manual or Automation,
    Sub Domain != TBD.
    """
    qa = 0
    not_qa = 0
    for r in data.get("defect_data", []):
        s = str(r.get("Sprint", "")).strip()
        if not s or s.startswith("Sprint "):
            continue
        t = str(r.get("Type Of Defect", "")).strip().lower()
        if t not in ("manual", "automation"):
            continue
        sub = str(r.get("Sub Domain", "")).strip()
        if not sub or sub == "TBD":
            continue
        count = _num(r.get("Prod Defect Count", 0))
        if count <= 0:
            continue
        rca = str(r.get("Prod Defect RCA", "") or "").strip().lower()
        if "qa miss" in rca:
            qa += count
        else:
            not_qa += count
    return qa, not_qa


def update_slide_15(slide, data):
    """Update Slide 15 (Production & Post-production Defect View)."""
    totals = _prod_defects_by_subdomain(data)
    grand = sum(totals.values())

    shape_map = {s.name: s for s in slide.shapes}

    # Chart 12 — Domain Wise Defects (proportions of grand total)
    if "Chart 12" in shape_map and shape_map["Chart 12"].has_chart:
        chart = shape_map["Chart 12"].chart
        try:
            series_name = chart.plots[0].series[0].name or "Domain Wise Defects"
        except Exception:
            series_name = "Domain Wise Defects"
        display_labels = [d for d, _ in SLIDE_15_DOMAIN_LABELS]
        values = []
        for _, sub_key in SLIDE_15_DOMAIN_LABELS:
            v = totals.get(sub_key, 0)
            values.append((v / grand) if grand > 0 else 0)
        update_chart_data(chart, display_labels, {series_name: values})

    # Chart 13 — QA Miss vs Not QA Miss distribution (by Prod Defect Count)
    if "Chart 13" in shape_map and shape_map["Chart 13"].has_chart:
        chart = shape_map["Chart 13"].chart
        try:
            series_name = chart.plots[0].series[0].name or "QA Miss Distribution"
        except Exception:
            series_name = "QA Miss Distribution"
        qa, not_qa = _qa_miss_split(data)
        tot = qa + not_qa
        if tot > 0:
            update_chart_data(chart, ["QA Miss", "Not QA Miss"],
                              {series_name: [qa / tot, not_qa / tot]})
        else:
            update_chart_data(chart, ["QA Miss", "Not QA Miss"],
                              {series_name: [0, 0]})

    # Title — replace month range portion dynamically
    months = _monthly_buckets_in_data(data)
    month_range = _format_month_range(months)
    if month_range:
        for s in slide.shapes:
            if s.has_text_frame and "PRODUCTION" in s.text_frame.text.upper() and "DEFECT" in s.text_frame.text.upper():
                _set_text_preserve_format(s, f"PRODUCTION & Post-production Defect View – {month_range}")
                break


def update_slide_16(slide, data):
    _update_rad_slide(slide, data, RAD_SLIDE_16_PANELS)


def update_slide_17(slide, data):
    _update_rad_slide(slide, data, RAD_SLIDE_17_PANELS)


# ==============================================================================
# SLIDES 18-19: Primary Release Automation (ReleaseDayTestCaseSheet)
# ==============================================================================

# Per-panel config for the Primary Release Automation slides:
#   (label_shape, chart_shape, sub_domain_key, display_label)
RELEASE_SLIDE_18_PANELS = [
    ("TextBox 26", "Chart 4",  "ALM-Technology",                 "ALM-Technology"),
    ("TextBox 27", "Chart 10", "Corporate Systems(TFG)",         "Corporate Systems(TFG)"),
    ("TextBox 28", "Chart 15", "Custody, Clearing & Settlement", "Custody, Clearing & Settlement"),
    ("TextBox 29", "Chart 16", "Data",                           "Data"),
    ("TextBox 30", "Chart 18", "Infosec",                        "Infosec"),
    ("TextBox 31", "Chart 24", "Practice Management",            "Practice Management"),
]
RELEASE_SLIDE_19_PANELS = [
    ("TextBox 26", "Chart 4",  "SRC",                 "SRC"),
    ("TextBox 27", "Chart 10", "Service and Support", "Service and Support"),
    ("TextBox 28", "Chart 15", "Technology",          "Technology"),
    ("TextBox 29", "Chart 16", "Trading",             "Trading"),
]


def _latest_release_month(data):
    """Return the most-recent 'Release Month' bucket in ReleaseDayTestCaseSheet
    that has at least one row with a non-blank, non-TBD Sub Domain.

    Values are formatted <MonthName>'YY (e.g. \"April'26\"). Sort key is
    (year, month-index). Months whose rows all have blank/TBD Sub Domain are
    skipped — otherwise every panel would collapse to zero the moment a new
    Release Month is added before its Sub Domain cells are filled in.
    Returns None if no month qualifies.
    """
    def _key(bucket):
        try:
            name, yr = bucket.split("'")
            return (int(yr), list(_MONTH_SHORT.keys()).index(name.strip().lower()))
        except Exception:
            return (-1, -1)
    months_with_subdomain = set()
    for r in data.get("release_day", []):
        m = str(r.get("Release Month", "")).strip()
        sub = str(r.get("Sub Domain", "")).strip()
        if m and sub and sub != "TBD":
            months_with_subdomain.add(m)
    if not months_with_subdomain:
        return None
    return max(months_with_subdomain, key=_key)


def _compute_release_by_subdomain(data):
    """{sub_domain: {stories, feasible, automated}} from ReleaseDayTestCaseSheet.

    Sums across ALL Release Month rows per Sub Domain
    (skipping TBD / blank Sub Domain). No month filter.
    """
    result = {}
    for r in data.get("release_day", []):
        sub = str(r.get("Sub Domain", "")).strip()
        if not sub or sub == "TBD":
            continue
        if sub not in result:
            result[sub] = {"stories": 0, "feasible": 0, "automated": 0}
        result[sub]["stories"]   += _num(r.get("No Of Stories part of the Release", 0))
        result[sub]["feasible"]  += _num(r.get("No of Automation feasible TCs", 0))
        result[sub]["automated"] += _num(r.get("No of TCs Automated", 0))
    return result


def _format_release_label(display, pct):
    """Format the per-panel percentage label, matching the template's style."""
    if display is None:
        return None
    if pct <= 0:
        return f"{display} ~ 0% stories automated"
    if pct >= 99.5:
        return f"{display} ~ 100% stories automated"
    return f"{display} ~ {pct:.1f}% stories automated"


def _update_release_slide(slide, data, panels):
    """Populate one Primary Release Automation slide using its panel config.

    For each panel:
      - chart series = [Story Volume, Automation Feasible TCs, Automated TCs]
        (sums across all rows of the panel's Sub Domain in ReleaseDayTestCaseSheet)
      - label = '<Sub Domain> ~ <Automated/Feasible×100>% stories automated'
    """
    release_data = _compute_release_by_subdomain(data)
    shape_map = {s.name: s for s in slide.shapes}
    cats = ["Story Volume", "Automation Feasible TCs", "Automated TCs"]

    for label_name, chart_name, sub_key, display in panels:
        vals = release_data.get(sub_key, {"stories": 0, "feasible": 0, "automated": 0})

        # Chart
        if chart_name in shape_map and shape_map[chart_name].has_chart:
            chart = shape_map[chart_name].chart
            try:
                series_name = chart.plots[0].series[0].name or "Series 1"
            except Exception:
                series_name = "Series 1"
            update_chart_data(chart, cats, {series_name: [
                int(vals["stories"]),
                int(vals["feasible"]),
                int(vals["automated"]),
            ]})

        # Label (skip blank panels)
        if label_name in shape_map and display is not None:
            pct = round(vals["automated"] / vals["feasible"] * 100, 1) if vals["feasible"] else 0
            _set_text_preserve_format(shape_map[label_name], _format_release_label(display, pct))


def update_slide_18(slide, data):
    _update_release_slide(slide, data, RELEASE_SLIDE_18_PANELS)


def update_slide_19(slide, data):
    _update_release_slide(slide, data, RELEASE_SLIDE_19_PANELS)


# ==============================================================================
# SLIDES 20-21: Regression Automation Trend (TCsDetailsSheet)
# ==============================================================================

# Panel config: (chart_name, sub_domain_key). One panel per Sub Domain.
TREND_SLIDE_20_PANELS = [
    ("Chart 14", "ALM-Technology"),
    ("Chart 15", "Corporate Systems(TFG)"),
    ("Chart 20", "Custody, Clearing & Settlement"),
    ("Chart 21", "Data"),
    ("Chart 22", "Infosec"),
    ("Chart 23", "Practice Management"),
    ("Chart 24", "Service and Support"),
]
TREND_SLIDE_21_PANELS = [
    ("Chart 12", "SRC"),
    ("Chart 16", "Technology"),
    ("Chart 14", "Trading"),
]

_MONTH_ABBR_UPPER = ["JAN", "FEB", "MAR", "APR", "MAY", "JUN",
                     "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]


def _detect_latest_month_abbr(data):
    """Return uppercase 3-letter abbreviation of the latest month found in
    TCsDetailsSheet's 'Submitted On' column (e.g. 'MAR'). None if unparseable.
    """
    import re
    latest = (0, 0)
    for r in data.get("tcs_details", []):
        so = str(r.get("Submitted On", "")).strip()
        m = re.match(r"(\d{1,2})/(\d{1,2})/(\d{4})", so)  # M/D/YYYY
        if m:
            mo, yr = int(m.group(1)), int(m.group(3))
        else:
            m = re.match(r"(\d{4})-(\d{2})-(\d{2})", so)  # YYYY-MM-DD
            if not m:
                continue
            yr, mo = int(m.group(1)), int(m.group(2))
        if 1 <= mo <= 12 and (yr, mo) > latest:
            latest = (yr, mo)
    if latest == (0, 0):
        return None
    return _MONTH_ABBR_UPPER[latest[1] - 1]


def _build_team_subdomain_lookup(data):
    """{lower-case scrum team name: sub_domain} built from ProjectDetails."""
    lookup = {}
    for r in data.get("project_details", []):
        team = str(r.get("Scrum Team Name", "")).strip().lower()
        sub = str(r.get("Sub Domain", "")).strip()
        if team and sub and sub != "TBD":
            lookup[team] = sub
    return lookup


def _compute_regression_by_subdomain(data):
    """{sub_domain: {automatable, automated}} from TCsDetailsSheet.

    Each TCsDetailsSheet row is mapped to a Sub Domain via the Scrum Team Name
    looked up in ProjectDetails. Rows whose team is not in ProjectDetails (or
    whose ProjectDetails Sub Domain is blank / TBD) are skipped.
    """
    team_lookup = _build_team_subdomain_lookup(data)
    result = {}
    for r in data.get("tcs_details", []):
        team = str(r.get("Scrum Team Name", "")).strip().lower()
        sub = team_lookup.get(team)
        if not sub:
            continue
        if sub not in result:
            result[sub] = {"automatable": 0, "automated": 0}
        for p in ["P0", "P1", "P2", "P3", "P4"]:
            result[sub]["automatable"] += _num(r.get(f"Total TCs Feasible {p}", 0))
            result[sub]["automated"]   += _num(r.get(f"Total TCs Automated {p}", 0))
    return result


def _update_trend_slide(slide, data, panels):
    """Populate one Regression Automation Trend slide.

    For each panel, the chart's existing categories (e.g. ['JAN', 'FEB']) and
    historical series values are preserved. The current Excel's totals are
    written into the latest-month column:
      - If the latest month already exists as a category, that column is replaced.
      - Otherwise the latest month is appended as a NEW column.
    """
    totals = _compute_regression_by_subdomain(data)
    latest = _detect_latest_month_abbr(data) or ""
    shape_map = {s.name: s for s in slide.shapes}

    for chart_name, sub_key in panels:
        if chart_name not in shape_map or not shape_map[chart_name].has_chart:
            continue
        chart = shape_map[chart_name].chart
        vals = totals.get(sub_key, {"automatable": 0, "automated": 0})

        # Capture existing chart structure
        existing_cats = [str(c) for c in chart.plots[0].categories]
        try:
            series_list = list(chart.plots[0].series)
            automatable_name = (series_list[0].name if series_list else None) or "Automatable"
            automated_name = (series_list[1].name if len(series_list) > 1 else None) or "Automated"
            old_automatable = [(v if v is not None else 0) for v in (list(series_list[0].values) if series_list else [])]
            old_automated = [(v if v is not None else 0) for v in (list(series_list[1].values) if len(series_list) > 1 else [])]
        except Exception:
            automatable_name, automated_name = "Automatable", "Automated"
            old_automatable, old_automated = [], []
        # Pad to category length
        while len(old_automatable) < len(existing_cats):
            old_automatable.append(0)
        while len(old_automated) < len(existing_cats):
            old_automated.append(0)

        new_aut = int(vals["automatable"])
        new_done = int(vals["automated"])

        if not latest:
            # No latest month detected — leave chart untouched
            continue

        if latest in existing_cats:
            # The latest month already exists as a column — update ONLY its
            # values (older columns stay as preserved historical data).
            idx = existing_cats.index(latest)
            new_cats = list(existing_cats)
            new_automatable = list(old_automatable); new_automatable[idx] = new_aut
            new_automated   = list(old_automated);   new_automated[idx]   = new_done
        else:
            # Roll the trend forward: keep the last 2 historical columns and
            # append the latest month as the 3rd. Older columns are dropped.
            new_cats = (existing_cats + [latest])[-3:]
            new_automatable = (list(old_automatable) + [new_aut])[-3:]
            new_automated   = (list(old_automated)   + [new_done])[-3:]

        # Always trim to at most 3 months — the chart is a 3-month rolling view.
        if len(new_cats) > 3:
            new_cats = new_cats[-3:]
            new_automatable = new_automatable[-3:]
            new_automated   = new_automated[-3:]

        update_chart_data(chart, new_cats, {
            automatable_name: new_automatable,
            automated_name:   new_automated,
        })

        # After replace_data: the chart's series XML keeps any per-point
        # data-label overrides (c:dLbl idx="N") from the template. Those
        # overrides drive the on-bar number's font/color/position for the
        # points they cover. When we add a NEW data point (the latest month),
        # it has no per-point override and falls back to the series-level
        # defaults — which in this template use dark text. Against dark bars
        # that becomes invisible.
        # Fix: for each series, if there is at least one per-point dLbl,
        # clone it for every category index that doesn't yet have one, so
        # the new bar's label inherits the same formatting as the historical ones.
        _extend_trend_datalabels(chart, len(new_cats))


def _extend_trend_datalabels(chart, n_cats):
    """Ensure every category index 0..n_cats-1 has a c:dLbl override matching
    the styling of the existing per-point overrides on this series. Called
    after python-pptx's replace_data(), which does not touch dLbl elements.
    """
    from pptx.oxml.ns import qn
    from copy import deepcopy
    C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    for ser in chart.plots[0].series:
        ser_xml = ser._element
        dLbls = ser_xml.find(qn("c:dLbls"))
        if dLbls is None:
            continue
        existing_dLbl = dLbls.findall(qn("c:dLbl"))
        if not existing_dLbl:
            continue  # No per-point overrides at all — series-level covers all
        # Find which indices already have per-point overrides
        covered = set()
        for dl in existing_dLbl:
            idx_el = dl.find(qn("c:idx"))
            if idx_el is not None:
                try:
                    covered.add(int(idx_el.get("val")))
                except (TypeError, ValueError):
                    pass
        # Clone the first override as a template for any missing index
        template = existing_dLbl[0]
        for i in range(n_cats):
            if i in covered:
                continue
            clone = deepcopy(template)
            # Update its idx
            new_idx = clone.find(qn("c:idx"))
            if new_idx is not None:
                new_idx.set("val", str(i))
            # Strip any c:extLst inside the clone that carries a uniqueId —
            # we don't want two dLbl elements sharing the same GUID.
            for ext_lst in clone.findall(qn("c:extLst")):
                clone.remove(ext_lst)
            # Insert BEFORE the series-level trailing elements — the schema
            # wants c:dLbl children to appear before c:spPr / c:txPr / etc.
            # Simplest: insert right after the last existing c:dLbl.
            last_dLbl_pos = 0
            for j, child in enumerate(list(dLbls)):
                if child.tag == qn("c:dLbl"):
                    last_dLbl_pos = j + 1
            dLbls.insert(last_dLbl_pos, clone)


def update_slide_20(slide, data):
    _update_trend_slide(slide, data, TREND_SLIDE_20_PANELS)


def update_slide_21(slide, data):
    _update_trend_slide(slide, data, TREND_SLIDE_21_PANELS)


# ==============================================================================
# SLIDES 6-9: Presence of Cognizant QE Team (ProjectDetails)
# ==============================================================================

# Pentagon shape → TextBox shape name (by vertical row position)
# Pentagon 11 (row 1) ↔ TextBox 21
# Pentagon 12 (row 2) ↔ TextBox 17
# Pentagon 13 (row 3) ↔ TextBox 19

PRESENCE_SLIDE_DOMAINS = {
    6: [
        ("Arrow: Pentagon 11", "TextBox 21", "Chart 10", "ALM-Technology"),
        ("Arrow: Pentagon 12", "TextBox 17", "Chart 9", "Corporate Systems(TFG)"),
        ("Arrow: Pentagon 13", "TextBox 19", "Chart 6", "Custody, Clearing & Settlement"),
    ],
    7: [
        ("Arrow: Pentagon 11", "TextBox 21", "Chart 14", "Data"),
        ("Arrow: Pentagon 12", "TextBox 17", "Chart 10", "Infosec"),
        ("Arrow: Pentagon 13", "TextBox 19", "Chart 9", "Practice Management"),
    ],
    8: [
        ("Arrow: Pentagon 11", "TextBox 21", "Chart 10", "Service and Support"),
        ("Arrow: Pentagon 12", "TextBox 17", "Chart 9", "SRC"),
        ("Arrow: Pentagon 13", "TextBox 19", "Chart 6", "Technology"),
    ],
    9: [
        (None, "TextBox 21", "Chart 9", "Trading"),
    ],
}

# Sub-domain display labels (how they appear in the PPT pentagon)
SUBDOMAIN_DISPLAY = {
    "ALM-Technology": "ALM - Technology",
    "Corporate Systems(TFG)": "Corporate Systems (TFG)",
    "Custody, Clearing & Settlement": "Custody, Clearing & Settlement",
    "Data": "Data",
    "Infosec": "Infosec",
    "Practice Management": "Practice Management",
    "Service and Support": "Service and Support",
    "SRC": "SRC",
    "Technology": "Technology",
    "Trading": "Trading",
}


def _build_team_lookup(data):
    """Build {sub_domain: {'SVT': [teams], 'Non SVT': [teams]}} from ProjectDetails."""
    lookup = {}
    for r in data.get("project_details", []):
        sub = str(r.get("Sub Domain", "")).strip()
        team = str(r.get("Scrum Team Name", "")).strip()
        ttype = str(r.get("Team Type", "")).strip()
        if not team:
            continue
        # Skip rows where Sub Domain is blank / TBD — they should not be
        # grouped into any sub-domain chart or text box.
        if not sub or sub in ("TBD", "nan", "None"):
            continue
        key = sub
        if key not in lookup:
            lookup[key] = {"SVT": [], "Non SVT": []}
        if ttype == "SVT":
            lookup[key]["SVT"].append(team)
        elif ttype == "Non SVT":
            lookup[key]["Non SVT"].append(team)
        # Skip TBD or any other Team Type (don't count in charts/text boxes)
    return lookup


def _format_teams_text(svt_teams, nonsvt_teams):
    """Format team list into readable text for the slide text box."""
    lines = []
    if svt_teams:
        lines.append("SVT Teams:")
        for t in svt_teams:
            lines.append(f"  \u2022 {t}")
    if nonsvt_teams:
        if lines:
            lines.append("")
        lines.append("Non SVT Teams:")
        for t in nonsvt_teams:
            lines.append(f"  \u2022 {t}")
    return "\n".join(lines) if lines else "No teams assigned"


def _capture_run_format(run):
    """Capture font formatting from a run to replicate later."""
    font = run.font
    fmt = {
        "name": font.name,
        "size": font.size,
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color_rgb": None,
        "color_theme": None,
    }
    try:
        if font.color and font.color.type is not None:
            try:
                fmt["color_rgb"] = font.color.rgb
            except AttributeError:
                pass
            try:
                fmt["color_theme"] = font.color.theme_color
            except AttributeError:
                pass
    except Exception:
        pass
    return fmt


def _apply_run_format(run, fmt):
    """Apply previously captured formatting to a run."""
    font = run.font
    if fmt.get("name"):
        font.name = fmt["name"]
    if fmt.get("size"):
        font.size = fmt["size"]
    if fmt.get("bold") is not None:
        font.bold = fmt["bold"]
    if fmt.get("italic") is not None:
        font.italic = fmt["italic"]
    if fmt.get("underline") is not None:
        font.underline = fmt["underline"]
    if fmt.get("color_rgb"):
        try:
            font.color.rgb = fmt["color_rgb"]
        except Exception:
            pass


def _replace_textframe_content(shape, lines, header_fmt=None, body_fmt=None):
    """Replace the text frame content with given lines preserving formatting.

    lines: list of (text, is_header) tuples where is_header=True gets header_fmt.
    If header_fmt/body_fmt are None, they are captured from the existing content.
    """
    tf = shape.text_frame

    # Capture existing formatting from first paragraph's first run if not provided
    if header_fmt is None or body_fmt is None:
        captured = None
        for p in tf.paragraphs:
            if p.runs:
                captured = _capture_run_format(p.runs[0])
                break
        if header_fmt is None:
            header_fmt = dict(captured) if captured else {}
            header_fmt["bold"] = True
        if body_fmt is None:
            body_fmt = dict(captured) if captured else {}
            body_fmt["bold"] = False  # Body team names should not be bold

    # Remove all paragraphs except the first (by XML manipulation)
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    paragraphs = txBody.findall(qn("a:p"))
    for p in paragraphs[1:]:
        txBody.remove(p)

    # Clear the first paragraph's content but keep its properties
    first_p = tf.paragraphs[0]
    for run in list(first_p.runs):
        run._r.getparent().remove(run._r)

    # Add lines back
    first = True
    for line_text, is_header in lines:
        if first:
            p = first_p
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line_text
        _apply_run_format(run, header_fmt if is_header else body_fmt)


def _update_presence_slide(slide, data, slide_num):
    """Update domain pentagon labels and SVT/Non SVT count charts for slides 6-9.

    NOTE: The 'Teams and Key Feature Implemented' text boxes (TextBox 17/19/21) are
    intentionally NOT updated — their feature description content is static and
    maintained manually by the author.
    """
    team_lookup = _build_team_lookup(data)
    domain_rows = PRESENCE_SLIDE_DOMAINS.get(slide_num, [])

    # Build shape name → shape map
    shape_map = {s.name: s for s in slide.shapes}

    for pentagon_name, textbox_name, chart_name, sub_domain in domain_rows:
        # 1. Update pentagon domain label
        if pentagon_name and pentagon_name in shape_map:
            display = SUBDOMAIN_DISPLAY.get(sub_domain, sub_domain)
            _set_text_preserve_format(shape_map[pentagon_name], display)

        # Look up teams for this sub-domain (with fuzzy fallback)
        teams = team_lookup.get(sub_domain, {"SVT": [], "Non SVT": []})
        if not teams["SVT"] and not teams["Non SVT"]:
            for key in team_lookup:
                if sub_domain.lower() in key.lower() or key.lower() in sub_domain.lower():
                    teams = team_lookup[key]
                    break

        # 2. Update SVT/Non SVT count chart
        if chart_name and chart_name in shape_map and shape_map[chart_name].has_chart:
            chart = shape_map[chart_name].chart
            svt_count = len(teams["SVT"])
            nonsvt_count = len(teams["Non SVT"])
            # Preserve existing category (e.g. year "2026") from the template
            try:
                existing_cats = [str(c) for c in chart.plots[0].categories]
                category = existing_cats[0] if existing_cats else "2026"
            except Exception:
                category = "2026"
            update_chart_data(chart, [category], {
                "SVT": [svt_count if svt_count > 0 else None],
                "Non SVT": [nonsvt_count if nonsvt_count > 0 else None],
            })


def update_slide_6(slide, data):
    _update_presence_slide(slide, data, 6)


def update_slide_7(slide, data):
    _update_presence_slide(slide, data, 7)


def update_slide_8(slide, data):
    _update_presence_slide(slide, data, 8)


def update_slide_9(slide, data):
    _update_presence_slide(slide, data, 9)


# ==============================================================================
# MAIN: Apply all updates to a presentation
# ==============================================================================

SLIDE_UPDATERS = {
    6: update_slide_6,
    7: update_slide_7,
    8: update_slide_8,
    9: update_slide_9,
    10: update_slide_10,
    11: update_slide_11,
    12: update_slide_12,
    13: update_slide_13,
    14: update_slide_14,
    15: update_slide_15,
    16: update_slide_16,
    17: update_slide_17,
    18: update_slide_18,
    19: update_slide_19,
    20: update_slide_20,
    21: update_slide_21,
}


def update_presentation(prs, data):
    """Apply all data updates to the presentation slides."""
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num in SLIDE_UPDATERS:
            try:
                SLIDE_UPDATERS[slide_num](slide, data)
            except Exception as e:
                print(f"Warning: Failed to update slide {slide_num}: {e}")
